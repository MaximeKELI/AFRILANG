#!/usr/bin/env python3
"""PowerPoint webinaire ISF Capsule IT — AFRILANG : UI/UX + animations."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT / "docs" / "webinars" / "AFRILANG_Capsule_IT_ISF_16_08_2026.pptx"

# Palette
NAVY = RGBColor(0x0A, 0x1A, 0x32)
NAVY_DEEP = RGBColor(0x06, 0x12, 0x24)
CARD = RGBColor(0x12, 0x2A, 0x48)
CARD2 = RGBColor(0x16, 0x32, 0x52)
ORANGE = RGBColor(0xF0, 0x7A, 0x1A)
ORANGE_SOFT = RGBColor(0xFF, 0xB0, 0x6B)
GREEN = RGBColor(0x2E, 0xC0, 0x5A)
GREEN_SOFT = RGBColor(0x7A, 0xE8, 0xA0)
CYAN = RGBColor(0x3A, 0xB4, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xE6, 0xEE, 0xF8)
MUTED = RGBColor(0x9A, 0xAE, 0xC4)
LINE = RGBColor(0x2A, 0x4A, 0x6E)

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

_anim_id = 100


def next_id() -> int:
    global _anim_id
    _anim_id += 1
    return _anim_id


def solid(shape, color: RGBColor, line: RGBColor | None = None, line_pt: float = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None or line_pt <= 0:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)


def send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bg(slide, top_accent=True):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    solid(bg, NAVY_DEEP)
    send_to_back(slide, bg)
    # Soft panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    solid(panel, NAVY)
    send_to_back(slide, panel)
    # Left accent rail
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    solid(rail, ORANGE)
    if top_accent:
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        )
        solid(top, ORANGE)


def footer(slide, page: int, total: int):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35)
    )
    solid(bar, RGBColor(0x08, 0x16, 0x2A))
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.045)
    )
    solid(accent, ORANGE)
    left = slide.shapes.add_textbox(Inches(0.45), Inches(7.22), Inches(10.2), Inches(0.25))
    p = left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "ISF  ·  Capsule IT S2E3  ·  AFRILANG — Coder l'Afrique de demain"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED
    r.font.name = "Calibri"
    right = slide.shapes.add_textbox(Inches(11.4), Inches(7.22), Inches(1.6), Inches(0.25))
    p2 = right.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page:02d}  /  {total:02d}"
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = ORANGE_SOFT
    r2.font.name = "Calibri"


def run(p, text, size=18, color=WHITE, bold=False, name="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = name
    return r


def card(slide, left, top, width, height, title, lines, accent=ORANGE, icon: str | None = None):
    """Carte UI : un seul shape (fond + texte) pour une animation propre."""
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.05),
        Inches(top + 0.06),
        Inches(width),
        Inches(height),
    )
    solid(shadow, RGBColor(0x04, 0x0C, 0x18))

    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    solid(body, CARD, accent, 1.75)

    # Accents barre haute via fine shape (non animée séparément — OK)
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left + 0.02), Inches(top + 0.02), Inches(width - 0.04), Inches(0.07)
    )
    solid(strip, accent)

    tf = body.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        body.text_frame.paragraphs[0].space_before = Pt(14)
    except Exception:
        pass
    title_text = f"{icon}  {title}" if icon else title
    p = tf.paragraphs[0]
    p.clear() if hasattr(p, "clear") else None
    # reset paragraph
    if p.runs:
        for r in list(p.runs):
            r.text = ""
    run(p, title_text, size=15, color=accent, bold=True)
    for line in lines:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run(p2, line, size=13, color=OFF)
    return body


def code_block(slide, left, top, width, height, code: str):
    """Returns the chrome shape (for animation targeting)."""
    chrome = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    solid(chrome, RGBColor(0x07, 0x12, 0x20), GREEN, 1.25)
    for i, c in enumerate(
        [RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]
    ):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + 0.2 + i * 0.28),
            Inches(top + 0.14),
            Inches(0.16),
            Inches(0.16),
        )
        solid(dot, c)
    label = slide.shapes.add_textbox(Inches(left + 1.2), Inches(top + 0.1), Inches(4), Inches(0.25))
    run(label.text_frame.paragraphs[0], "exemple.afr", size=11, color=MUTED)
    box = slide.shapes.add_textbox(
        Inches(left + 0.25),
        Inches(top + 0.45),
        Inches(width - 0.5),
        Inches(height - 0.6),
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run(p, line if line else " ", size=13, color=GREEN_SOFT, name="Consolas")
    return chrome


def heading(slide, title: str, subtitle: str | None = None, y=0.28):
    """Returns [title_box] for animation."""
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(y), Inches(2.35), Inches(0.32)
    )
    solid(pill, CARD2, ORANGE, 1.25)
    pt = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.02), Inches(2.25), Inches(0.28))
    p = pt.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "CAPSULE IT  ·  ISF", size=10, color=ORANGE_SOFT, bold=True)

    box = slide.shapes.add_textbox(Inches(0.45), Inches(y + 0.4), Inches(12.4), Inches(1.05))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, title, size=30, color=WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run(p2, subtitle, size=15, color=ORANGE_SOFT, bold=False)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y + 1.25), Inches(1.6), Inches(0.05)
    )
    solid(line, GREEN)
    return [box]


def bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        run(p, f"▸  {item}", size=size, color=OFF)
    return box



def shape_id(shape) -> str:
    return str(shape.shape_id)


def add_slide_transition(slide, kind: str = "fade", speed: str = "med"):
    """Add a slide transition (fade / push / wipe)."""
    cSld = slide._element.find(qn("p:cSld"))
    # Remove existing transition
    for child in list(slide._element):
        if child.tag == qn("p:transition"):
            slide._element.remove(child)
    tr = etree.SubElement(slide._element, qn("p:transition"))
    tr.set("spd", speed)
    tr.set("advClick", "1")
    if kind == "push":
        etree.SubElement(tr, qn("p:push")).set("dir", "l")
    elif kind == "wipe":
        etree.SubElement(tr, qn("p:wipe")).set("dir", "r")
    else:
        etree.SubElement(tr, qn("p:fade"))
    # Ensure transition is after cSld
    if cSld is not None:
        slide._element.remove(tr)
        cSld_index = list(slide._element).index(cSld)
        slide._element.insert(cSld_index + 1, tr)


def add_entrance_animations(slide, shapes, effect: str = "fade", stagger_ms: int = 180):
    """Attach sequential entrance animations to shapes (click to advance each)."""
    if not shapes:
        return
    # Flatten nested lists and filter None / duplicates
    flat = []
    stack = list(shapes)
    while stack:
        item = stack.pop(0)
        if item is None:
            continue
        if isinstance(item, (list, tuple)):
            stack = list(item) + stack
            continue
        flat.append(item)

    seen = set()
    ordered = []
    for sh in flat:
        try:
            sid = shape_id(sh)
        except Exception:
            continue
        if sid in seen:
            continue
        seen.add(sid)
        ordered.append(sh)
    if not ordered:
        return

    # Build timing XML
    timing = etree.Element(qn("p:timing"))
    tnLst = etree.SubElement(timing, qn("p:tnLst"))
    root_par = etree.SubElement(tnLst, qn("p:par"))
    root_cTn = etree.SubElement(root_par, qn("p:cTn"))
    root_cTn.set("id", str(next_id()))
    root_cTn.set("dur", "indefinite")
    root_cTn.set("restart", "never")
    root_cTn.set("nodeType", "tmRoot")
    root_children = etree.SubElement(root_cTn, qn("p:childTnLst"))

    seq = etree.SubElement(root_children, qn("p:seq"))
    seq.set("concurrent", "1")
    seq.set("nextAc", "seek")
    seq_cTn = etree.SubElement(seq, qn("p:cTn"))
    seq_cTn.set("id", str(next_id()))
    seq_cTn.set("dur", "indefinite")
    seq_cTn.set("nodeType", "mainSeq")
    seq_children = etree.SubElement(seq_cTn, qn("p:childTnLst"))

    # presetID: 10=fade, 2=fly, 22=wipe
    preset_map = {"fade": ("10", "0"), "fly": ("2", "4"), "wipe": ("22", "4")}
    preset_id, subtype = preset_map.get(effect, ("10", "0"))

    for i, sh in enumerate(ordered):
        sid = shape_id(sh)
        delay = "0" if i == 0 else str(stagger_ms)

        outer = etree.SubElement(seq_children, qn("p:par"))
        outer_cTn = etree.SubElement(outer, qn("p:cTn"))
        outer_cTn.set("id", str(next_id()))
        outer_cTn.set("fill", "hold")
        st = etree.SubElement(outer_cTn, qn("p:stCondLst"))
        cond = etree.SubElement(st, qn("p:cond"))
        # First on click/with previous click root; subsequent after previous
        if i == 0:
            cond.set("delay", "0")
        else:
            cond.set("delay", "0")
            # after previous via seq order — use clickEffect chain
        outer_children = etree.SubElement(outer_cTn, qn("p:childTnLst"))

        mid = etree.SubElement(outer_children, qn("p:par"))
        mid_cTn = etree.SubElement(mid, qn("p:cTn"))
        mid_cTn.set("id", str(next_id()))
        mid_cTn.set("fill", "hold")
        mid_st = etree.SubElement(mid_cTn, qn("p:stCondLst"))
        mid_cond = etree.SubElement(mid_st, qn("p:cond"))
        mid_cond.set("delay", delay if i else "0")
        mid_children = etree.SubElement(mid_cTn, qn("p:childTnLst"))

        inner = etree.SubElement(mid_children, qn("p:par"))
        inner_cTn = etree.SubElement(inner, qn("p:cTn"))
        inner_cTn.set("id", str(next_id()))
        inner_cTn.set("presetID", preset_id)
        inner_cTn.set("presetClass", "entr")
        inner_cTn.set("presetSubtype", subtype)
        inner_cTn.set("fill", "hold")
        inner_cTn.set("grpId", "0")
        inner_cTn.set("nodeType", "clickEffect")
        inner_st = etree.SubElement(inner_cTn, qn("p:stCondLst"))
        inner_cond = etree.SubElement(inner_st, qn("p:cond"))
        inner_cond.set("delay", "0")
        inner_children = etree.SubElement(inner_cTn, qn("p:childTnLst"))

        # Always use fade entrance (reliable across PowerPoint / Impress)
        anim_effect = etree.SubElement(inner_children, qn("p:animEffect"))
        anim_effect.set("transition", "in")
        anim_effect.set("filter", "fade")
        cBhvr = etree.SubElement(anim_effect, qn("p:cBhvr"))
        cTn = etree.SubElement(cBhvr, qn("p:cTn"))
        cTn.set("id", str(next_id()))
        cTn.set("dur", "450")
        tgt = etree.SubElement(cBhvr, qn("p:tgtEl"))
        spTgt = etree.SubElement(tgt, qn("p:spTgt"))
        spTgt.set("spid", sid)

        # set visibility visible at start of effect
        set_el = etree.Element(qn("p:set"))
        set_cBhvr = etree.SubElement(set_el, qn("p:cBhvr"))
        set_cTn = etree.SubElement(set_cBhvr, qn("p:cTn"))
        set_cTn.set("id", str(next_id()))
        set_cTn.set("dur", "1")
        set_st = etree.SubElement(set_cTn, qn("p:stCondLst"))
        set_cond = etree.SubElement(set_st, qn("p:cond"))
        set_cond.set("delay", "0")
        set_tgt = etree.SubElement(set_cBhvr, qn("p:tgtEl"))
        set_sp = etree.SubElement(set_tgt, qn("p:spTgt"))
        set_sp.set("spid", sid)
        attr = etree.SubElement(set_cBhvr, qn("p:attrNameLst"))
        an = etree.SubElement(attr, qn("p:attrName"))
        an.text = "style.visibility"
        to = etree.SubElement(set_el, qn("p:to"))
        val = etree.SubElement(to, qn("p:strVal"))
        val.set("val", "visible")
        inner_children.insert(0, set_el)

    # prevCondLst / nextCondLst for seq
    prev = etree.SubElement(seq, qn("p:prevCondLst"))
    prev_cond = etree.SubElement(prev, qn("p:cond"))
    prev_cond.set("evt", "onPrev")
    prev_cond.set("delay", "0")
    prev_tgt = etree.SubElement(prev_cond, qn("p:tgtEl"))
    etree.SubElement(prev_tgt, qn("p:sldTgt"))
    nxt = etree.SubElement(seq, qn("p:nextCondLst"))
    next_cond = etree.SubElement(nxt, qn("p:cond"))
    next_cond.set("evt", "onNext")
    next_cond.set("delay", "0")
    next_tgt = etree.SubElement(next_cond, qn("p:tgtEl"))
    etree.SubElement(next_tgt, qn("p:sldTgt"))

    # Replace / append timing on slide
    for child in list(slide._element):
        if child.tag == qn("p:timing"):
            slide._element.remove(child)
    slide._element.append(timing)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    return slide


def build():
    global _anim_id
    _anim_id = 100

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    anim_plan: list[tuple] = []  # (slide, shapes, effect, transition)

    def register(slide, shapes, effect="fade", transition="fade"):
        # Copy now — callers may reuse/mutate the same list name later.
        anim_plan.append((slide, list(shapes), effect, transition))

    logo = ASSETS / "afrilang-africa-clean.png"

    # ── 1 Titre ──
    s = new_slide(prs)
    anim = []
    if logo.exists():
        pic = s.shapes.add_picture(str(logo), Inches(0.7), Inches(1.55), height=Inches(2.4))
        anim.append(pic)
    hero = s.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(9.2), Inches(3.8))
    tf = hero.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "AFRILANG", size=52, color=WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run(p, "Coder l'Afrique de demain", size=28, color=ORANGE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    run(
        p,
        "Un langage pensé pour l'inclusion, l'innovation\net la souveraineté numérique africaine.",
        size=17,
        color=OFF,
    )
    anim.append(hero)
    meta = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.35), Inches(8.8), Inches(1.2)
    )
    solid(meta, CARD, ORANGE, 1.25)
    anim.append(meta)
    mt = s.shapes.add_textbox(Inches(3.7), Inches(5.5), Inches(8.4), Inches(1.0))
    tf = mt.text_frame
    run(tf.paragraphs[0], "Capsule IT — Saison 2 · Épisode 3", size=15, color=ORANGE_SOFT, bold=True)
    p = tf.add_paragraph()
    run(p, "Dimanche 16 août 2026  ·  19h00  ·  Google Meet", size=14, color=OFF)
    p = tf.add_paragraph()
    run(p, "Organisé par Informaticiens sans frontière (ISF)", size=13, color=MUTED)
    anim.append(mt)
    register(s, anim, "fade", "fade")

    # ── 2 Bienvenue ──
    s = new_slide(prs)
    h = heading(s, "Bienvenue à Capsule IT", "Informaticiens sans frontière — Togo")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Merci d'être présents pour ce webinaire en ligne.",
            "ISF : « L'informatique au cœur du développement ».",
            "Capsule IT démystifie les technologies et inspire l'action.",
            "Objectif : comprendre AFRILANG, sa vision, et comment démarrer.",
            "Format : présentation claire → démos → questions / réponses.",
            "Notez vos questions — on les traitera en fin de session.",
        ],
        size=17,
    )
    register(s, h + [b], "fade", "push")

    # ── 3 Agenda ──
    s = new_slide(prs)
    h = heading(s, "Au programme ce soir", "Parcours structuré (~60–90 min)")
    anim = list(h)
    agenda = [
        ("01", "Contexte & enjeux", "Souveraineté et inclusion", GREEN),
        ("02", "Intervenant", "Maxime Dzidula KELI · SIG", ORANGE),
        ("03", "AFRILANG", "Philosophie, syntaxe, archi", CYAN),
        ("04", "Cas d'usage", "Éducation, SIG, innovation", GREEN),
        ("05", "Communauté", "Roadmap & contribution", ORANGE),
        ("06", "Q & R", "Échanges ouverts", CYAN),
    ]
    for i, (num, title, sub, accent) in enumerate(agenda):
        col, row = i % 3, i // 3
        anim.append(card(
            s,
            0.45 + col * 4.2,
            1.85 + row * 2.35,
            4.0,
            2.15,
            f"{num}  {title}",
            [sub],
            accent,
            icon=None,
        ))
    register(s, anim, "fade", "fade")

    # ── 4 Objectifs ──
    s = new_slide(prs)
    h = heading(s, "Objectifs pédagogiques", "Ce que vous emportez ce soir")
    anim = list(h)
    anim.append(card(s, 0.45, 1.85, 4.0, 4.7, "Comprendre", [
        "Pourquoi un langage africain ?",
        "Inclusion ≠ simplification",
        "Souveraineté = maîtrise",
        "Place d'AFRILANG",
    ], GREEN, "◎"))
    anim.append(card(s, 4.65, 1.85, 4.0, 4.7, "Savoir faire", [
        "Lire la syntaxe naturelle",
        "Lancer un 1er programme",
        "Explorer le playground",
        "Relier à vos projets (SIG…)",
    ], ORANGE, "⚡"))
    anim.append(card(s, 8.85, 1.85, 4.0, 4.7, "Agir", [
        "Rejoindre la communauté",
        "Proposer des exemples",
        "Former autour de vous",
        "Contribuer docs / code",
    ], CYAN, "➜"))
    register(s, anim, "fade", "wipe")

    # ── 5 ISF ──
    s = new_slide(prs)
    h = heading(s, "Informaticiens sans frontière", "L'informatique au service du développement")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Association engagée pour démocratiser le numérique et former les talents.",
            "Capsule IT : webinaire accessible, régulier, ouvert à tous les niveaux.",
            "Passerelle entre technique, éducation et impact social.",
            "Contact : Contact@isftogo.com  ·  +228 90 94 30 67  ·  +228 98 65 90 00",
            "Site : http://isftogo.com",
            "Ce soir : Saison 2 · Épisode 3 — focus AFRILANG.",
        ],
        size=17,
    )
    register(s, h + [b], "fade", "fade")

    # ── 6 Intervenant ──
    s = new_slide(prs)
    h = heading(s, "Intervenant", "Maxime Dzidula KELI")
    anim = list(h)
    anim.append(card(s, 0.45, 1.85, 6.05, 4.7, "Profil", [
        "Ingénieur géoinformaticien.",
        "Solutions intégrant :",
        "  • Systèmes d'information géographique",
        "  • Technologies informatiques",
        "",
        "Pont entre territoire et logiciel —",
        "un axe stratégique pour l'Afrique.",
    ], ORANGE, "👤"))
    anim.append(card(s, 6.75, 1.85, 6.05, 4.7, "Pourquoi ce sujet ?", [
        "Le SIG structure agriculture,",
        "urbanisme, santé, logistique…",
        "",
        "Un langage accessible accélère :",
        "  • formation des équipes locales",
        "  • prototypage d'outils métier",
        "  • appropriation des données",
        "",
        "AFRILANG + géomatique =",
        "souveraineté code + territoire.",
    ], GREEN, "🗺"))
    register(s, anim, "fade", "push")

    # ── 7 Défi ──
    s = new_slide(prs)
    h = heading(s, "Le défi : coder pour l'Afrique, depuis l'Afrique")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Des millions de jeunes motivés — mais des barrières linguistiques et cognitives.",
            "Les langages dominants freinent souvent l'entrée en programmation.",
            "Forte dépendance à des outils conçus hors du continent.",
            "Peu de chaînes de compilation maîtrisées localement de bout en bout.",
            "Besoin d'outils pédagogiques adaptés au contexte éducatif africain.",
            "Question : former massivement sans perdre en exigence technique ?",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "fade")

    # ── 8 Souveraineté ──
    s = new_slide(prs)
    h = heading(s, "Souveraineté numérique : 4 leviers")
    anim = list(h)
    for i, (t, lines, accent) in enumerate(
        [
            ("Données", ["Qui stocke ?", "Qui analyse ?", "Qui décide ?"], GREEN),
            ("Logiciels", ["Qui écrit ?", "Qui audite ?", "Qui maintient ?"], ORANGE),
            ("Compétences", ["Qui forme ?", "Qui certifie ?", "Qui emploie ?"], CYAN),
            ("Infrastructures", ["Cloud local ?", "Open source ?", "Interop ?"], GREEN),
        ]
    ):
        anim.append(card(s, 0.4 + i * 3.2, 1.9, 3.05, 4.5, t, lines, accent))
    register(s, anim, "fade", "wipe")

    # ── 9 Inclusion ──
    s = new_slide(prs)
    h = heading(s, "Inclusion : baisser la barrière, pas le plafond")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Inclusion ≠ langage « jouet » : vraie puissance (OOP, types, modules…).",
            "Syntaxe naturelle : le code se rapproche de la pensée humaine.",
            "Idéal collèges, lycées, bootcamps, universités, autodidactes.",
            "Focus d'abord sur les algorithmes, puis la montée en compétence.",
            "Generics, patterns, async, GUI, packages… le plafond reste haut.",
            "AFRILANG : accessibilité à l'entrée + sérieux à la sortie.",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "fade")

    # ── 10 Pourquoi ──
    s = new_slide(prs)
    h = heading(s, "Pourquoi AFRILANG plutôt qu'un framework éducatif ?")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 6.05, 4.6, "Un vrai langage", [
        "Grammaire complète + compilateur",
        "stdlib, packages, tests, playground",
        "Transpile vers C++17 → perf & écosystème",
        "Cibles : natif, WASM, JS / Canvas",
        "GUI : std/ui, game2d, game3d",
    ], GREEN, "◆"))
    anim.append(card(s, 6.75, 1.9, 6.05, 4.6, "Une identité", [
        "L'Afrique au centre du geste de coder",
        "Narratif clair pour écoles & États",
        "Projet communautaire ancré localement",
        "Docs bilingues FR / EN",
        "Slogan : Coder l'Afrique",
    ], ORANGE, "✦"))
    register(s, anim, "fade", "push")

    # ── 11 Définition ──
    s = new_slide(prs)
    h = heading(s, "AFRILANG en une phrase")
    quote = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.3), Inches(11.9), Inches(3.4)
    )
    solid(quote, CARD, GREEN, 2)
    qt = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(11.1), Inches(2.8))
    tf = qt.text_frame
    tf.word_wrap = True
    run(
        tf.paragraphs[0],
        "Un langage orienté objet à syntaxe naturelle\n"
        "qui transpile vers C++17 — conçu pour apprendre,\n"
        "prototyper et déployer des logiciels depuis l'Afrique.",
        size=24,
        color=WHITE,
        bold=True,
    )
    p = tf.add_paragraph()
    p.space_before = Pt(18)
    run(p, "afrilang.allcoders.tech   ·   Coder l'Afrique", size=16, color=ORANGE_SOFT)
    register(s, h + [quote, qt], "fade", "fade")

    # ── 12 Philosophie ──
    s = new_slide(prs)
    h = heading(s, "Trois piliers", "Lisibilité · Puissance · Souveraineté")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 4.0, 4.6, "1. Lisibilité", [
        'say "Bonjour"',
        "if score is greater than 10…",
        "while window is open do…",
        "",
        "Le code se lit comme une explication.",
    ], GREEN, "①"))
    anim.append(card(s, 4.65, 1.9, 4.0, 4.6, "2. Puissance", [
        "Classes, interfaces, enums",
        "Génériques + where",
        "Modules, packages, macros",
        "Async / générateurs",
        "GUI 2D / 3D + stdlib",
    ], ORANGE, "②"))
    anim.append(card(s, 8.85, 1.9, 4.0, 4.6, "3. Souveraineté", [
        "Compilateur maîtrisé",
        "Chaîne native / WASM / web",
        "Docs bilingues",
        "Formation & communauté",
        "Vision continentale",
    ], CYAN, "③"))
    register(s, anim, "fade", "wipe")

    # ── 13 Hello ──
    s = new_slide(prs)
    h = heading(s, "Premier contact : Hello AFRILANG")
    anim = list(h)
    anim.append(code_block(
        s,
        0.45,
        1.85,
        6.4,
        4.7,
        '''say "Bonjour depuis AFRILANG!"

repeat 3 times
    say "Hello"
end

create score = 85
if score is greater than 79 then
    say "Bravo : {score}"
else
    say "Continue"
end''',
    ))
    anim.append(card(s, 7.15, 1.85, 5.6, 4.7, "Repères", [
        "say → afficher",
        "create / set → variables",
        "repeat / while / for each",
        "Conditions en anglais naturel",
        'Interpolation : "{score}"',
        "Commentaires : // …",
    ], ORANGE, "💡"))
    register(s, anim, "fade", "fade")

    # ── 14 Types ──
    s = new_slide(prs)
    h = heading(s, "Système de types", "Clarté pour apprendre, rigueur pour grandir")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Types de base : number, text, bool, listes, maps…",
            "Optionnels (text?) et Result pour gérer absence / erreur.",
            "Enums avec payloads : Status.Error with \"message\".",
            "Génériques : function identity<T>(x T) returns T",
            "Contraintes : where T is number",
            "Le compilateur aide tôt — moins de surprises à l'exécution.",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "push")

    # ── 15 OOP ──
    s = new_slide(prs)
    h = heading(s, "Orienté objet, sans jargon inutile")
    anim = list(h)
    anim.append(code_block(
        s,
        0.45,
        1.85,
        6.5,
        4.7,
        '''class Greeter
    function hello()
        say "Hello from a class!"
    end
end

create g = new Greeter
g.hello()''',
    ))
    anim.append(card(s, 7.2, 1.85, 5.55, 4.7, "Ce que ça ouvre", [
        "Classes & constructeurs",
        "Interfaces (contrats)",
        "Héritage (extends)",
        "Modélisation métier claire",
        "Projets scolaires structurés",
        "Apps réelles, pas jouets",
    ], GREEN, "🧩"))
    register(s, anim, "fade", "fade")

    # ── 16 Architecture ──
    s = new_slide(prs)
    h = heading(s, "Architecture du compilateur", "De .afr à l'exécutable")
    anim = list(h)
    steps = [
        ("Source", ".afr", GREEN),
        ("Lexer", "tokens", ORANGE),
        ("Parser", "AST", CYAN),
        ("Sémantique", "types", GREEN),
        ("CodeGen", "C++17", ORANGE),
        ("Backend", "natif/WASM/JS", CYAN),
    ]
    for i, (t, sub, accent) in enumerate(steps):
        anim.append(card(s, 0.35 + i * 2.15, 2.15, 2.05, 2.6, t, [sub], accent))
    note = s.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.2), Inches(1.3))
    tf = note.text_frame
    tf.word_wrap = True
    run(
        tf.paragraphs[0],
        "Pipeline : .afr → imports → Lexer → Parser → AST → Sémantique → CodeGen → g++ / em++ / JS",
        size=14,
        color=MUTED,
    )
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run(p, "Un langage source — plusieurs cibles selon le besoin.", size=14, color=OFF)
    anim.append(note)
    register(s, anim, "fade", "wipe")

    # ── 17 Backends ──
    s = new_slide(prs)
    h = heading(s, "Trois façons d'exécuter")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 4.0, 4.6, "Natif (C++)", [
        "afrilang run fichier.afr",
        "Performance maximale",
        "SDL pour GUI desktop",
        "Apps & jeux locaux",
    ], GREEN, "🖥"))
    anim.append(card(s, 4.65, 1.9, 4.0, 4.6, "Playground web", [
        "Browser / Canvas / WebGL",
        "Apprendre sans installer",
        "UI + game2d + game3d",
        "afrilang.allcoders.tech",
    ], ORANGE, "🌐"))
    anim.append(card(s, 8.85, 1.9, 4.0, 4.6, "WASM", [
        "Exécution navigateur",
        "Modules portables",
        "Intégration sites / apps",
        "Chaîne Emscripten",
    ], CYAN, "📦"))
    register(s, anim, "fade", "fade")

    # ── 18 Stdlib ──
    s = new_slide(prs)
    h = heading(s, "Bibliothèque standard & packages")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Modules std : math, collections, json, http, fs, crypto, datetime…",
            "Graphisme : std/ui, std/game2d, std/game3d.",
            "État de jeu : gamestate (menus, pause, game over).",
            "Packages partageables (afrilang pkg …).",
            "Documentation web bilingue FR / EN.",
            "Exemples : hello, OOP, snake, démos 3D, finance…",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "push")

    # ── 19 GUI ──
    s = new_slide(prs)
    h = heading(s, "Interfaces graphiques", "Du bouton simple au monde 3D")
    anim = list(h)
    anim.append(code_block(
        s,
        0.45,
        1.85,
        6.5,
        4.7,
        '''open window titled "AFRILANG GUI"
    with width 640, height 480

while window is open do
    clear background color 28, 28, 36
    draw text "Bonjour !" at 60, 80 size 32
    if button "Quitter" at 220, 350
        width 200 height 50 is clicked then
        close window
    end
    show frame
end''',
    ))
    anim.append(card(s, 7.2, 1.85, 5.55, 4.7, "Modules GUI", [
        "std/ui — fenêtres & widgets",
        "std/game2d — grilles, sprites",
        "std/game3d — OpenGL / WebGL",
        "Playground Canvas navigateur",
        "Pédagogie visuelle forte",
        "Projets motivants étudiants",
    ], CYAN, "🎮"))
    register(s, anim, "fade", "fade")

    # ── 20 Playground ──
    s = new_slide(prs)
    h = heading(s, "Le playground : coder en 30 secondes", "https://afrilang.allcoders.tech")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Éditeur en ligne + exemples prêts à lancer.",
            "Run natif (serveur) pour programmes texte.",
            "Browser : JS/Canvas pour GUI sans installer SDL.",
            "Workspace multi-fichiers dans le navigateur.",
            "Parfait ateliers ISF, clubs, cours à distance.",
            "Confidentialité & consentement intégrés au site.",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "wipe")

    # ── 21 Éducation ──
    s = new_slide(prs)
    h = heading(s, "Cas d'usage : éducation & formation")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 4.0, 4.6, "Collège / Lycée", [
        "Premiers algorithmes",
        "Boucles & conditions",
        "Mini-jeux 2D",
        "Clubs informatiques",
    ], GREEN, "🏫"))
    anim.append(card(s, 4.65, 1.9, 4.0, 4.6, "Université / Bootcamp", [
        "OOP & architecture",
        "Tests & packages",
        "Compilation & backends",
        "Projets GUI / WASM",
    ], ORANGE, "🎓"))
    anim.append(card(s, 8.85, 1.9, 4.0, 4.6, "Associations (ISF…)", [
        "Ateliers Capsule IT",
        "Hackathons locaux",
        "Mentorat pair-à-pair",
        "Contenus FR accessibles",
    ], CYAN, "🤝"))
    register(s, anim, "fade", "fade")

    # ── 22 SIG ──
    s = new_slide(prs)
    h = heading(s, "Cas d'usage : géomatique & territoire", "Pont avec l'expertise de Maxime Dzidula KELI")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Le SIG structure la décision publique et privée.",
            "Besoin de scripts / outils métiers adaptés localement.",
            "AFRILANG : langage d'apprentissage avant Python/C++ SIG.",
            "Prototypes : calculs, imports, petites UI de saisie.",
            "Vision : packages formats géo, APIs carto, ETL légers.",
            "Former des géomaticiens qui comprennent aussi le logiciel.",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "push")

    # ── 23 Innovation ──
    s = new_slide(prs)
    h = heading(s, "Cas d'usage : innovation locale")
    anim = list(h)
    for i, (t, lines, accent) in enumerate(
        [
            ("AgriTech", ["Suivi parcelles", "Alertes saisonnières", "Coopératives"], GREEN),
            ("EdTech", ["Exercices interactifs", "Jeux éducatifs", "Évaluations"], ORANGE),
            ("FinTech", ["Micro-gestion", "Budgets coopératifs", "Éducation financière"], CYAN),
            ("CivicTech", ["Open data local", "Cartes citoyennes", "Services mairie"], GREEN),
        ]
    ):
        anim.append(card(s, 0.4 + i * 3.2, 1.9, 3.05, 4.5, t, lines, accent))
    register(s, anim, "fade", "wipe")

    # ── 24 Démo ──
    s = new_slide(prs)
    h = heading(s, "Fil conducteur de démonstration (live)")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "1) Hello + conditions dans le playground (2 min).",
            "2) Petite classe Greeter — OOP accessible (3 min).",
            "3) gui_demo : fenêtre + bouton Quitter (3 min).",
            "4) Aperçu snake / game2d : motivation visuelle (3 min).",
            "5) Mention game3d_demo : même langage, monde 3D (2 min).",
            "6) Où trouver docs, exemples, téléchargements (2 min).",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "fade")

    # ── 25 Installer ──
    s = new_slide(prs)
    h = heading(s, "Démarrer en local", "Pour aller plus loin que le navigateur")
    anim = list(h)
    anim.append(code_block(
        s,
        0.5,
        1.85,
        12.3,
        4.7,
        '''# Construire
cmake -S . -B build && cmake --build build -j

# Lancer
./build/afrilang run examples/hello.afr
./build/afrilang run examples/gui_demo.afr

# Playground (sans install lourde)
# https://afrilang.allcoders.tech/playground/''',
    ))
    register(s, anim, "fade", "push")

    # ── 26 Communauté ──
    s = new_slide(prs)
    h = heading(s, "Communauté & contribution")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "Contribuer : exemples, traductions, docs, correctifs.",
            "Partager : ateliers ISF, clubs, tutoriels vidéo FR.",
            "Proposer : besoins métier (SIG, finance, éducation).",
            "Tester : remonter bugs playground & plateformes.",
            "Réseau : développeurs, enseignants, décideurs.",
            "Éthique : inclusion, ouverture, exigence technique.",
        ],
        size=16,
    )
    register(s, h + [b], "fade", "fade")

    # ── 27 Roadmap ──
    s = new_slide(prs)
    h = heading(s, "Vision & roadmap")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 4.0, 4.6, "Court terme", [
        "Playground GUI renforcé",
        "Docs & tutoriels",
        "Exemples métiers",
        "Ateliers Capsule IT",
    ], GREEN, "▶"))
    anim.append(card(s, 4.65, 1.9, 4.0, 4.6, "Moyen terme", [
        "Packages géo / data",
        "Meilleur support IDE",
        "Formations certifiantes",
        "Partenariats écoles",
    ], ORANGE, "▶▶"))
    anim.append(card(s, 8.85, 1.9, 4.0, 4.6, "Long terme", [
        "Écosystème continental",
        "Outils souverains métiers",
        "Communauté multi-pays",
        "Standard éducatif local",
    ], CYAN, "★"))
    register(s, anim, "fade", "wipe")

    # ── 28 Takeaways ──
    s = new_slide(prs)
    h = heading(s, "À retenir ce soir", "5 messages")
    b = bullets(
        s,
        0.55,
        1.85,
        12.2,
        4.8,
        [
            "1. L'Afrique a besoin d'outils qu'elle comprend, enseigne et maîtrise.",
            "2. AFRILANG allie syntaxe naturelle et backend sérieux (C++ / WASM / web).",
            "3. Inclusion et excellence technique ne s'opposent pas.",
            "4. Le playground abaisse la barrière ; le compilateur lève le plafond.",
            "5. La souveraineté numérique commence par des compétences locales.",
        ],
        size=17,
    )
    register(s, h + [b], "fade", "fade")

    # ── 29 CTA ──
    s = new_slide(prs)
    h = heading(s, "Appel à l'action", "Que faire dès demain ?")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 6.05, 4.6, "Apprenants", [
        "1. Ouvrir le playground",
        "2. Lancer hello + gui_demo",
        "3. Modifier un exemple",
        "4. Partager à un ami / club",
        "5. Poser une question à ISF",
    ], GREEN, "🚀"))
    anim.append(card(s, 6.75, 1.9, 6.05, 4.6, "Structures", [
        "1. Planifier un atelier AFRILANG",
        "2. Module d'initiation",
        "3. Relier à un besoin SIG / métier",
        "4. Soutenir la documentation FR",
        "5. Co-construire des packages",
    ], ORANGE, "🏛"))
    register(s, anim, "fade", "push")

    # ── 30 Q&R ──
    s = new_slide(prs)
    h = heading(s, "Questions / Réponses")
    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.4), Inches(10.9), Inches(3.5)
    )
    solid(panel, CARD, ORANGE, 2)
    qt = s.shapes.add_textbox(Inches(1.6), Inches(2.9), Inches(10.1), Inches(2.6))
    tf = qt.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "Merci pour votre attention.", size=34, color=WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    run(
        p,
        "Questions sur le langage, l'installation, la pédagogie,\n"
        "le SIG, la feuille de route, la contribution…",
        size=18,
        color=OFF,
    )
    register(s, h + [panel, qt], "fade", "fade")

    # ── 31 Contacts ──
    s = new_slide(prs)
    h = heading(s, "Contacts & ressources")
    anim = list(h)
    anim.append(card(s, 0.45, 1.9, 6.05, 4.6, "Informaticiens sans frontière", [
        "Email : Contact@isftogo.com",
        "Tél. : +228 90 94 30 67",
        "Tél. : +228 98 65 90 00",
        "Web : http://isftogo.com",
        "",
        "Capsule IT — S2E3",
        "Dim 16/08/2026 · 19h00 · Meet",
    ], ORANGE, "☎"))
    anim.append(card(s, 6.75, 1.9, 6.05, 4.6, "AFRILANG", [
        "Site : afrilang.allcoders.tech",
        "Playground en ligne",
        "Documentation FR / EN",
        "Exemples & stdlib",
        "",
        "Slogan : Coder l'Afrique",
        "Intervenant : Maxime Dzidula KELI",
    ], GREEN, "🌍"))
    register(s, anim, "fade", "wipe")

    # ── 32 Remerciements ──
    s = new_slide(prs)
    h = heading(s, "Remerciements")
    b = bullets(
        s,
        0.55,
        1.95,
        12.2,
        4.5,
        [
            "Merci à Informaticiens sans frontière pour Capsule IT.",
            "Merci à Maxime Dzidula KELI pour le partage d'expertise.",
            "Merci à toutes les personnes connectées ce soir.",
            "Merci aux contributeurs qui font vivre AFRILANG.",
            "Ensemble, codons l'Afrique de demain — avec exigence et inclusion.",
        ],
        size=18,
    )
    register(s, h + [b], "fade", "fade")

    # ── 33 Final ──
    s = new_slide(prs)
    anim = []
    if logo.exists():
        pic = s.shapes.add_picture(str(logo), Inches(5.55), Inches(1.15), height=Inches(2.1))
        anim.append(pic)
    box = s.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "AFRILANG", size=46, color=WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run(p2, "Coder l'Afrique de demain", size=26, color=ORANGE, bold=True)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(18)
    run(p3, "ISF  ·  Capsule IT S2E3  ·  16 août 2026", size=15, color=MUTED)
    anim.append(box)
    register(s, anim, "fade", "fade")

    # Apply animations + transitions FIRST (before footers)
    for slide, shapes, effect, transition in anim_plan:
        add_slide_transition(slide, transition)
        key = []
        for sh in shapes:
            if sh is None:
                continue
            if isinstance(sh, (list, tuple)):
                key.extend([x for x in sh if x is not None and not isinstance(x, (list, tuple))])
            else:
                key.append(sh)
        add_entrance_animations(slide, key[:10], effect=effect, stagger_ms=100)

    # Footers last (not animated)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"Slides: {total}")
    print("UI: accent rail, cards, code chrome, footers")
    print("UX: fade/push/wipe transitions + click-to-reveal entrance animations")


if __name__ == "__main__":
    build()
